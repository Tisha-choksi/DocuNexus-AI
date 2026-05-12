from pathlib import Path

from pptx import Presentation


def extract_pptx(path: Path) -> tuple[list[dict], dict[str, str]]:
    presentation = Presentation(path)
    pages = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        pages.append({"page_number": index, "text": "\n".join(lines)})

    props = presentation.core_properties
    metadata = {
        "author": props.author or "",
        "title": props.title or "",
        "subject": props.subject or "",
        "created": props.created.isoformat() if props.created else "",
        "modified": props.modified.isoformat() if props.modified else "",
        "slide_count": str(len(presentation.slides)),
    }
    metadata = {key: value for key, value in metadata.items() if value}
    return pages, metadata

